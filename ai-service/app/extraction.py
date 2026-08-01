from pathlib import Path
import csv
import re


ASSET_PATTERN = re.compile(r"\b[A-Z]{1,4}[- ]?\d{2,4}[A-Z]?\b")
DATE_PATTERN = re.compile(
    r"\b(?:\d{1,2}[-/](?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[-/]\d{2,4}"
    r"|\d{1,2}[-/]\d{1,2}[-/]\d{2,4}|\d{4}-\d{2}-\d{2})\b",
    re.IGNORECASE,
)
MEASUREMENT_PATTERN = re.compile(
    r"\b\d+(?:\.\d+)?\s?(?:bar|°?c|rpm|mm/s|psi|kpa|mpa|%)\b", re.IGNORECASE
)
FAILURE_TERMS = {
    "bearing wear", "bearing failure", "seal leakage", "cavitation",
    "overheating", "high vibration", "lubrication contamination",
    "misalignment", "corrosion", "pressure loss", "motor trip",
}
ACTION_TERMS = {
    "replaced", "inspected", "lubricated", "aligned", "cleaned", "tightened",
    "overhauled", "calibrated", "tested", "flushed",
}


def extract_text(path: Path) -> tuple[str, list[dict]]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _pdf(path)
    if suffix == ".docx":
        return _docx(path)
    if suffix == ".csv":
        return _csv(path)
    if suffix in {".jpg", ".jpeg", ".png", ".tiff", ".bmp"}:
        return _image(path)
    if suffix in {".xlsx", ".xls"}:
        return _xlsx(path)
    if suffix == ".pptx":
        return _pptx(path)
    if suffix in {".html", ".htm"}:
        return _html(path)
    if suffix in {".eml", ".msg"}:
        return _eml(path)
    text = path.read_text(encoding="utf-8", errors="ignore")
    return text, [{"page": 1, "text": text}]


def _pdf(path: Path) -> tuple[str, list[dict]]:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pages.append({"page": index, "text": text})
    full_text = "\n\n".join(item["text"] for item in pages)
    if len(full_text.strip()) < 30:
        ocr_text = _ocr_pdf(path)
        if ocr_text:
            pages = [{"page": index + 1, "text": text} for index, text in enumerate(ocr_text)]
            full_text = "\n\n".join(ocr_text)
    return full_text, pages


def _ocr_pdf(path: Path) -> list[str]:
    try:
        import pytesseract
        from pdf2image import convert_from_path
        from .config import settings
        if settings.tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd
        return [pytesseract.image_to_string(image) for image in convert_from_path(str(path), dpi=180)]
    except Exception:
        return []


def _docx(path: Path) -> tuple[str, list[dict]]:
    from docx import Document
    document = Document(str(path))
    text = "\n".join(p.text for p in document.paragraphs if p.text.strip())
    return text, [{"page": 1, "text": text}]


def _csv(path: Path) -> tuple[str, list[dict]]:
    rows = []
    with path.open("r", encoding="utf-8-sig", errors="ignore", newline="") as handle:
        for row in csv.reader(handle):
            rows.append(" | ".join(row))
    text = "\n".join(rows)
    return text, [{"page": 1, "text": text}]


def _image(path: Path) -> tuple[str, list[dict]]:
    try:
        import pytesseract
        from PIL import Image
        from .config import settings
        if settings.tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd
        text = pytesseract.image_to_string(Image.open(str(path)))
        return text, [{"page": 1, "text": text}]
    except Exception as e:
        return f"Image OCR extraction failed: {e}", [{"page": 1, "text": ""}]


def _xlsx(path: Path) -> tuple[str, list[dict]]:
    try:
        import pandas as pd
        df = pd.read_excel(str(path))
        rows = []
        for _, row in df.iterrows():
            rows.append(" | ".join(f"{col}: {val}" for col, val in row.items() if pd.notna(val)))
        text = "\n".join(rows)
        return text, [{"page": 1, "text": text}]
    except Exception as e:
        return f"Excel parsing failed: {e}", [{"page": 1, "text": ""}]


def _pptx(path: Path) -> tuple[str, list[dict]]:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        pages = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        if p.text.strip():
                            slide_text.append(p.text.strip())
            pages.append({"page": idx, "text": "\n".join(slide_text)})
        full_text = "\n\n".join(item["text"] for item in pages)
        return full_text, pages
    except Exception as e:
        return f"PPTX parsing failed: {e}", [{"page": 1, "text": ""}]


def _html(path: Path) -> tuple[str, list[dict]]:
    try:
        from bs4 import BeautifulSoup
        text = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser").get_text(separator="\n")
        return text, [{"page": 1, "text": text}]
    except Exception:
        text = path.read_text(encoding="utf-8", errors="ignore")
        return text, [{"page": 1, "text": text}]


def _eml(path: Path) -> tuple[str, list[dict]]:
    try:
        import email
        from email import policy
        msg = email.message_from_string(path.read_text(encoding="utf-8", errors="ignore"), policy=policy.default)
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    body += part.get_content()
        else:
            body = msg.get_content()
        text = f"Subject: {msg.get('subject', '')}\nFrom: {msg.get('from', '')}\nDate: {msg.get('date', '')}\n\n{body}"
        return text, [{"page": 1, "text": text}]
    except Exception:
        text = path.read_text(encoding="utf-8", errors="ignore")
        return text, [{"page": 1, "text": text}]


def extract_entities(text: str) -> dict:
    lowered = text.lower()
    asset_tags = {
        normalize_asset(value) for value in ASSET_PATTERN.findall(text)
        if normalize_asset(value).split("-", 1)[0] not in {"IR", "SOP", "OEM", "ISO", "API"}
    }
    return {
        "asset_tags": sorted(asset_tags),
        "dates": sorted(set(DATE_PATTERN.findall(text)))[:50],
        "measurements": sorted(set(MEASUREMENT_PATTERN.findall(text)), key=str.lower)[:50],
        "failures": sorted(term for term in FAILURE_TERMS if term in lowered),
        "actions": sorted(term for term in ACTION_TERMS if term in lowered),
    }


def normalize_asset(value: str) -> str:
    return value.upper().replace(" ", "-")


def classify_document(name: str, text: str) -> str:
    value = f"{name} {text[:1000]}".lower()
    for needle, label in [
        ("inspection", "Inspection Report"),
        ("incident", "Incident Report"),
        ("maintenance", "Maintenance Record"),
        ("operating procedure", "Standard Operating Procedure"),
        ("sop", "Standard Operating Procedure"),
        ("manual", "Equipment Manual"),
        ("safety", "Safety Document"),
        ("reading", "Operating Readings"),
    ]:
        if needle in value:
            return label
    return "Industrial Document"
