import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Function to add a slide with title and content
def add_slide(prs, title_text, content_text):
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x4d, 0x40) # Dark Teal
    
    tf = content.text_frame
    tf.text = content_text
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.space_after = Pt(14)
        
    return slide

# Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "IndusMind AI"
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x4d, 0x40)
subtitle.text = "Industrial Asset Intelligence & RCA Copilot\n\nSpeaker 1: Introduction\nSpeaker 2: Frontend\nSpeaker 3: AI Brain\nSpeaker 4: Backend"

# Slide 2: Problem Statement & Solution
add_slide(prs, "The Problem & Our Solution", 
          "• Engineers waste hours digging through fragmented maintenance logs and OEM manuals.\n"
          "• Solution: An AI-powered knowledge layer.\n"
          "• Ask natural language questions (e.g., 'Why did Pump P-101 fail?')\n"
          "• Get instant answers with exact citations and evidence.\n"
          "• Tech Stack: React (Frontend), Spring Boot (Backend), FastAPI + LLaMA-3 (AI).")

# Slide 3: Frontend Experience
add_slide(prs, "The Frontend Experience", 
          "• Built with React 19, TypeScript, and Vite.\n"
          "• Fast, responsive, and intuitive single-page application.\n"
          "• Secure login via JWT authentication.\n"
          "• Features:\n"
          "   - Command Center Dashboard\n"
          "   - AI Copilot Interface\n"
          "   - Asset 360° Timeline View\n"
          "• Not a generic chatbot: specialized interface for industrial use cases.")

# Slide 4: The AI Brain (RAG & LoRA)
add_slide(prs, "The AI Brain: Privacy & Accuracy", 
          "• 100% On-Premise: No data sent to the cloud.\n"
          "• Retrieval-Augmented Generation (RAG):\n"
          "   - Embedded ChromaDB vector database.\n"
          "   - Exact page-level source citations.\n"
          "• Fine-Tuned Local LLM:\n"
          "   - LLaMA-3 8B model fine-tuned using LoRA.\n"
          "   - 4-bit quantization for fast local GPU execution.\n"
          "   - Optimized for industrial engineering terminology.")

# Slide 5: Backend & Conclusion
add_slide(prs, "Backend Architecture & Conclusion", 
          "• Core API: Java 21 & Spring Boot.\n"
          "• Strict Security & Role-Based Access Control (RBAC).\n"
          "• PostgreSQL Database management.\n"
          "• Conclusion:\n"
          "   - Turns unstructured data into actionable intelligence.\n"
          "   - Highly specialized, production-ready AI tool.\n"
          "   - Guarantees data privacy for heavy industry.")

# Save the presentation
prs.save('IndusMind_Final_Presentation.pptx')
print("Successfully generated IndusMind_Final_Presentation.pptx!")
